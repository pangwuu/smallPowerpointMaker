"""
Module for holding functions that create slides
"""
import os
from random import choice
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from helpers import scripts_folder
from ccli import find_ccli
# from tkinter import filedialog, Tk
from functools import cache
from deep_translator import GoogleTranslator
from ai_translate import  translate_with_gemini
import tempfile

class Song:
    '''
    Used to hold the data of a song

    Title - song title
    ccli - ccli message
    lyrics - a 2D array of song lyrics
    '''

    def __init__(self, title, ccli, lyrics):
        self.title = title
        self.ccli = ccli
        self.lyrics = lyrics

class MyPresentation:
    '''
    Used to return presentation values and their associated font size.
    '''
    def __init__(self, presentation, font_size):
        self.presentation = presentation
        self.font_size = font_size

def create_blank_slide(prs):
    '''
    Adds a blank slide using the powerpoint oject

    Each template must have the 7th template to be a blank slide
    '''
    layout = prs.slide_layouts[6]  # Blank slide layout. 
    return prs.slides.add_slide(layout)

def add_text_to_slide(slide, text, prs, font_size, alignment=PP_ALIGN.CENTER, position_percent=0.35,
                     colour=None, bold=False, italic=False, underline=False):
    '''
    Adds text aligned horizontally in the middle of the page, with position_percent deciding the height of the text
    colour is a string corresponding to a hex colour without the leading hashtag. If set to None, it defaults to the
    font colour specified by the current theme
    '''
    width = prs.slide_width
    height = prs.slide_height
    left = int((width - width * 0.8) / 2)
    top = int(height * position_percent)
    text_box = slide.shapes.add_textbox(left, top, int(width * 0.8), int(height * 0.5))
    text_frame = text_box.text_frame
    text_frame.text = text
    
    # Format ALL paragraphs, not just the first one
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.alignment = alignment
        paragraph.font.bold = bold
        paragraph.font.italic = italic
        paragraph.font.underline = underline
        if colour:
            paragraph.font.color.rgb = RGBColor.from_string(colour)
    
    text_frame.word_wrap = True

def create_bulletin_slide(slide, prs, date, songs, verses, speaker="TBA", topic="TBA"):
    '''
    Creates a slide with contents similar to that of the BCCC English Service bulletin front page
    '''
    # Set the background to be solid white using a shape
    # This ensures themes that use shapes in their background are properly hidden
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # \x0b is a special line break character, equivalent to Shift + Enter in PowerPoint
    # This ensures any font attributes (e.g. colour, size, etc.) are applied to all lines in the text box
    ppt_text_break = "\x0b"
    sunday_date = date.split("_")
    sunday_date_string = f"{sunday_date[2]}.{sunday_date[1]}.{sunday_date[0]}"
    add_text_to_slide(slide, "Blacktown Chinese Christian Church", prs, 30, position_percent=0.05,
                      bold=True, colour="000000")
    add_text_to_slide(slide, f"Welcome to Our English Service{ppt_text_break}{sunday_date_string}", prs, 23, position_percent=0.15,
                      bold=True, colour="000000")
    # 7030A0 = Dark purple
    add_text_to_slide(slide, "English Service 9:45 – 11:00 am", prs, 20, position_percent=0.35,
                      colour="7030A0", bold=True, underline=True)
    # Note that 'Introduction and Opening Prayer' is currently omitted
    bulletin_summary = {
        "left": [
            f"Worship Songs:{ppt_text_break * (len(songs) - 1)}",
            f"Passage:{ppt_text_break * (len(verses) - 1)}",
            "Speaker:",
            "Topic:",
            "Response Song:",
            "Announcements and Closing Prayer",
            "Benediction",
        ],
        "right": [
            ppt_text_break.join(songs),
            ppt_text_break.join(verses),
            speaker,
            topic,
            "TBA",
            "",
            "",
        ],
    }
    # python-pptx currently doesn't support text boxes having modifiers such as boldness applied to selective parts
    # of the text, so the alternative is to have two text boxes placed at the same position but with different
    # alignments so that one of them can apply boldness
    add_text_to_slide(slide, ppt_text_break.join(bulletin_summary["left"]), prs, 20, position_percent=0.4,
                      alignment=PP_ALIGN.LEFT, bold=True, colour="000000")
    add_text_to_slide(slide, ppt_text_break.join(bulletin_summary["right"]), prs, 20, position_percent=0.4,
                      alignment=PP_ALIGN.CENTER, colour="000000")

def create_title_slide(title_text: str, subtitle_text: str, prs, title_size, default_body_size=8):
    '''
    Creates a custom slide with a title text and body text.
    Usually used for song titles with small descriptions
    '''
    blank_slide = create_blank_slide(prs)

    subtitle_text = subtitle_text.replace("\n", " ")

    add_text_to_slide(blank_slide, title_text, prs, title_size, position_percent=0.2)
    add_text_to_slide(blank_slide, subtitle_text, prs, default_body_size, position_percent=0.6)

    return prs

def add_ccli_slide(prs, ccli_info, font_size):
    if ccli_info is None:
        return prs
    
    ccli_info = ccli_info.replace("\n", " ")
    
    blank_slide = create_blank_slide(prs)
    add_text_to_slide(blank_slide, ccli_info, prs, font_size, position_percent=0.4)
    
    return prs

def create_title_and_text_slide(title_text, body_text, prs, title_size,
                                body_size, slide_number=0, total_slides=0,song_mode=False):
    '''
    Creates a normal slide with title text and smaller body text. Result is similar to the "Title with content" slide
    Usually used in song lyrics slides
    '''

    # Make a new blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Use the blank slide layout
    lyric_slide = prs.slides.add_slide(blank_slide_layout)

    # Calculate the center and top position for title text box
    title_width = prs.slide_width * 0.9
    title_height = prs.slide_height * 0.15

    title_left = (prs.slide_width - title_width) / 2
    title_top = Inches(0.3)  # Adjust vertically as needed

    # Add a text box for the title
    title_box = lyric_slide.shapes.add_textbox(left=title_left, top=title_top, width=title_width, height=title_height)
    title_frame = title_box.text_frame
    title_frame.text = title_text

    # Calculate the center and top position for body text box
    body_width = prs.slide_width * 0.9
    body_height = prs.slide_height * 0.8
    body_left = (prs.slide_width - body_width) / 2
    body_top = title_height + title_top  # Adjust vertically as needed

    # Add a text box for the body text
    body_box = lyric_slide.shapes.add_textbox(left=body_left, top=body_top, width=body_width, height=body_height)
    body_frame = body_box.text_frame
    body_frame.text = body_text

    # Enable text wrapping for both title and body text
    title_frame.word_wrap = True
    body_frame.word_wrap = True

    # Set font size and boldness for both title and body text
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(title_size)  # Adjust the font size as needed
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(body_size)  # Adjust the font size as needed
        paragraph.font.bold = False
        paragraph.alignment = PP_ALIGN.CENTER


    # Presents the little box on the bottom right to show if slides are changing
    # if song_mode:
    #     prs = make_slide_box(prs, lyric_slide, slide_number, total_slides)

    return prs

def create_text_slide(title_text, body_text, prs, title_size,
                                body_size, slide_number=0, total_slides=0,song_mode=False):
    '''
    Creates a normal slide with just smaller body text. Result is similar to the "Content" slide
    Usually used in song lyrics slides
    '''

    # Make a new blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Use the blank slide layout
    lyric_slide = prs.slides.add_slide(blank_slide_layout)

    # Calculate the center and top position for title text box
    title_width = prs.slide_width * 0.9
    title_height = prs.slide_height * 0.15

    title_left = (prs.slide_width - title_width) / 2

    # Calculate the center and top position for body text box
    body_width = prs.slide_width * 0.9
    body_height = prs.slide_height * 0.8
    body_left = (prs.slide_width - body_width) / 2
    body_top = title_height  # Adjust vertically as needed

    # Add a text box for the body text
    body_box = lyric_slide.shapes.add_textbox(left=body_left, top=body_top, width=body_width, height=body_height)
    body_frame = body_box.text_frame
    body_frame.text = body_text

    # Enable text wrapping for both title and body text
    body_frame.word_wrap = True

    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(body_size)  # Adjust the font size as needed
        paragraph.font.bold = False
        paragraph.alignment = PP_ALIGN.CENTER

    # Presents the little box on the bottom right to show if slides are changing
    # if song_mode:
    #     prs = make_slide_box(prs, lyric_slide, slide_number, total_slides)

    return prs

def make_slide_box(prs, lyric_slide, slide_number, total_slides):
    # Define dimensions and positions relative to slide width and height
    margin_right = prs.slide_width * 0.10
    margin_bottom = prs.slide_height * 0.10
    round_rect_width = prs.slide_width * 0.2
    round_rect_height = prs.slide_height * 0.1

    # Create round rectangle shape
    round_rect = lyric_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        prs.slide_width - margin_right - round_rect_width,
        prs.slide_height - margin_bottom - round_rect_height,
        round_rect_width,
        round_rect_height
    )

    # Set round rectangle fill color
    fill = round_rect.fill
    fill.solid()
    fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_2

    # Add text box with slide number and total slides on top of round rectangle
    textbox = lyric_slide.shapes.add_textbox(
        prs.slide_width - margin_right - round_rect_width,
        prs.slide_height - margin_bottom - round_rect_height,
        round_rect_width,
        round_rect_height
    )
    # Add the little box on the bottom
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"{slide_number + 1}/{total_slides}"

    # Center align text
    p.alignment = PP_ALIGN.CENTER

    # Change font properties if needed
    font = p.font
    
    font.bold = True
    return prs

def song_object_from_name(song_name: str, max_lines: int=4) -> Song:
    '''
    Reads the content from a song text file and returns a song object with filled data
    '''
    
    song_name = song_name.lower().title().strip()
    # Locate where the text file is stored
    lyrics_text_file = f"{scripts_folder}/../Songs/{song_name}/{song_name}_Lyrics.txt"

    new_song = Song('', '', '')

    with open(lyrics_text_file, encoding='utf-8') as file:

        # Extract all parts into a dictionary
        lines = file.readlines()

        lyrics = []

        line_number = 0

        current_section = None

        for line in lines:
            line_number += 1

            # First two lines are reserved for the song name and ccli description
            if line_number <= 2:
                continue

            line = line.strip()
            if line.strip().startswith("[") and line.strip().endswith("]"):
                # New section detected - append the current section to the lyrics.
                # Every section corresponds to a slide that will be shown
                current_section = [line[1:-1].strip(), '']
                lyrics.append(current_section)

            elif current_section is not None:
                # Check if the current lyrics exceed a maximum number of lines
                if len(current_section[1].split('\n')) > max_lines:
                    # If we exceed, create a new section with the same title
                    current_section = [current_section[0], '']
                    lyrics.append(current_section)

                # Append lyrics to the current section
                current_section[1] += line + '\n'
        
        # Create a new song object with the required data
        new_song = Song(lines[0], lines[1], lyrics)

    if new_song.title == '' and new_song.ccli == '' and new_song.lyrics == '':
        raise FileNotFoundError(f'The song of name {song_name} does not seem to exist. Check the Songs directory to see if it is there.')
    return new_song

def append_song_to_powerpoint(song_name, prs, title_size, font_size, max_lines=4):
    '''
    Will append a song's lyrics from a text file with all its lyrics to the current powerpoint. 
    Operates by selecting the required song from the Songs folder and using the contained .txt file to create a song with its associated parts
    
    The text file format MUST be, including the square brackets:
    Song name
    CCLI: [Number here]
    [Section title]
    Lyrics for the section
    [Section title]
    Lyrics for the section
    '''

    # Copy a song - I'm gonna use a specific one for now - we need to fix up a search algorithm for this that works
    try:
        new_song = song_object_from_name(song_name)
    except FileNotFoundError:
        print("Song not found - continuing without adding it")
        return

    # First slide of the song with title data and ccli data
    song_ccli_info = find_ccli(new_song.title.replace("\n", "").replace("(live)", "").strip().lower())
    prs = create_title_slide(new_song.title.strip().lower().title().replace(' (Live)', ''), song_ccli_info, prs, title_size)
    for slide_number, lyrics in enumerate(new_song.lyrics):
        # Insert a generic lyrics slide for each set of lyrics that exist
        prs = create_text_slide(lyrics[0], lyrics[1], prs,title_size, font_size, slide_number=slide_number, total_slides=len(new_song.lyrics), song_mode=True)
    
    return prs

def create_from_template(test_mode=False, select_template=False) -> Presentation:
    '''
    Creates a new PowerPoint based on existing templates in the Templates folder.
    Allows users to select a specific template if needed via a GUI file picker.
    '''
    directory_path = f"{scripts_folder}/../Templates"
    files_to_keep = [os.path.join(root, file) for root, _, files in os.walk(directory_path) for file in files if file.endswith(".pptx")]
    
    if not files_to_keep:
        raise FileNotFoundError("No PowerPoint templates found in the Templates folder.")
    
    # If in test mode, return a list of template presentations
    if test_mode:
        ppts_to_return = []
        for file in files_to_keep:
            prs = Presentation(file)
            basename = os.path.basename(file)
            size = 'small' if basename.startswith('small') else 'medium' if basename.startswith('med') else 'large'
            ppts_to_return.append(MyPresentation(prs, size))
        return ppts_to_return
    
    selected_template = choice(files_to_keep)
    
    prs = Presentation(selected_template)
    return prs, selected_template


def create_starting_slides(prs: Presentation, title_size: int, left_text_size: int) -> Presentation:
    # Bulletin slide
    # This is initially empty since it relies on data sourced from user input later on, but the slide is still created
    # at this point to ensure the final slide ordering is correct
    create_blank_slide(prs)

    # Title slide
    create_title_slide('BCCC english service', '', prs, title_size)

    # Please be quiet slide
    add_title_with_image_on_right(prs, "Please keep all phones on silent", 'Phone', left_text_size)
    
    # Worship slide
    # add_title_with_image_on_right(prs, 'Worship', 'Worship', left_text_size)
    return prs

def add_title_with_image_on_right(prs: Presentation, title_text: str, image_type: str, left_text_size: int) -> Presentation:
    '''
    Adds a slide with a title on the left and a large square image on the right.
    Make sure images are square. Non square images will be distorted into square ones and will look weird
    Image type corresponds to the folder the image will be randomly selected from, in "Images"
    '''
    margin_right = margin_top = prs.slide_height * 0.1

    # Create a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Slide layout index 6 represents a blank slide

    # Add title
    slide_title = slide.shapes.add_textbox(prs.slide_width * 0.05, margin_top, prs.slide_width * 0.4, prs.slide_height - 2 * margin_top)
    title_text_frame = slide_title.text_frame
    title_text_frame.word_wrap = True
    title_text_frame.auto_size = True
    p = title_text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(left_text_size)

    # Choose a file from the specified type
    directory_path = f"{scripts_folder}/../Images/{image_type.strip().lower().title()}"
    files = [f for f in os.listdir(directory_path) if os.path.isfile(os.path.join(directory_path, f))]
    
    if '.DS_Store' in files:
        files.remove('.DS_Store')
    
    if not files:
        return None  # No files in the directory

    # Choose a random image
    for i in range(100):
        random_image = choice(files)
        if random_image.endswith('png'):
            break

    image_path = os.path.join(directory_path, random_image)

    image_width = image_height = prs.slide_height - 2 * margin_top

    # Add image on the right side
    slide.shapes.add_picture(image_path, prs.slide_width - image_width - margin_right, margin_top, width=image_width,
                              height=image_height)

    return prs

# This code is clunky but does the job
def create_offering_slide(prs: Presentation, tithing_heading_size: int, tithing_body_size: int) -> Presentation:
    '''
    Creates a slide for tithing details.
    Creates 5 roundrects to store the tithing details with text boxes with required information on them.
    Icons are found within the tithing folder
    '''
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add a text box for the title
    title_height = Inches(1)
    title_top = (prs.slide_height - title_height) / 2

    left_textbox = prs.slide_width*0.075
    top_textbox = title_top
    width_textbox = prs.slide_width*0.3
    height_textbox = Inches(1)

    textbox = slide.shapes.add_textbox(left_textbox, top_textbox - height_textbox/2, width_textbox, height_textbox)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()

    # Set the text and font size for the title in the text box
    p.text = 'Offering'
    
    p.font.size = Pt(tithing_heading_size)

    # Optional: Adjust other text box properties if needed
    text_frame.word_wrap = True
    text_frame.auto_size = True
    text_frame.margin_left = 0  # Adjust left margin as needed

    # Boxes margin
    box_height = prs.slide_height/6
    box_spacing = box_height*0.15

    # Calculate the total height occupied by the 5 boxes and the spacing
    total_height = 5 * box_height + 4 * box_spacing

    # Calculate the top position for the first box to be 0.4 inch from the top
    top_first_box = (prs.slide_height - total_height) / 2

    # Insert the 5 shapes which form the tithing details
    for i in range(5):
        width_rect = 0.65 * prs.slide_width # 35% of the slide
        height_rect = box_height

        left_rect = prs.slide_width - (top_first_box + width_rect) # Distance from left
        top_rect = top_first_box + i * (box_spacing + box_height) # Distance from top

        # Add rounded rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left_rect, top_rect, width_rect, height_rect
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
        fill.fore_color.brightness = -0.15
        shape.line.fill.background()

        # Add a text box on top of the rounded rectangle
        left_textbox = left_rect + width_rect/7
        top_textbox = top_rect  # Adjust top position based on your design
        # if i != 1 and i != 2:
        #     top_textbox = top_rect + box_height/9  # Adjust top position based on your design

        width_textbox = width_rect
        height_textbox = height_rect

        textbox = slide.shapes.add_textbox(left_textbox, top_textbox, width_textbox - width_rect/5, height_textbox)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()


        image_width = image_height = box_height*0.8
        p.font.size = Pt(tithing_body_size)
        p.font.name = "Gill Sans MT"

        # Text differs depending on which box you're writing in
        if i == 0:
            p.text = "Account name: Blacktown Chinese Christian Church"
            slide.shapes.add_picture(f"{scripts_folder}/../Images/Tithing/church.png", left_rect + box_height*0.1, top_rect + box_height*0.1, image_width, image_height)
        elif i == 1:
            p.text = "Account number: 4216 50263"
            slide.shapes.add_picture(f"{scripts_folder}/../Images/Tithing/account_number.png", left_rect  + box_height*0.1, top_rect + box_height*0.1, image_width, image_height)
        elif i == 2:
            p.text = "BSB: 112 - 879"
            slide.shapes.add_picture(f"{scripts_folder}/../Images/Tithing/bsb.png", left_rect  +  box_height*0.1, top_rect + box_height*0.1, image_width, image_height)
        elif i == 3:
            p.text = "Please put in \"offering\" as the reference"
            slide.shapes.add_picture(f"{scripts_folder}/../Images/Tithing/hands.png", left_rect +  box_height*0.1, top_rect + box_height*0.1, image_width, image_height)
        elif i == 4:         
            p.text = "The offering box is available at the back of the hall"
            slide.shapes.add_picture(f"{scripts_folder}/../Images/Tithing/box.png", left_rect +  box_height*0.1, top_rect + box_height*0.1, image_width, image_height)

        text_frame.word_wrap = True  # Enable word wrapping
        text_frame.auto_size = True  # Enable autofit

    return prs

@cache
def translate_text(text: str, language: str = 'mandarin') -> str:
    '''
    Takes in a string and translates it into a language given by the language parameter (default - Mandarin simplified)
    '''

    # mapping
    top_languages = {
        "Mandarin Chinese": "zh-CN",
        "Spanish": "es",
        "Hindi": "hi",
        "Arabic": "ar",
        "Bengali": "bn",
        "Portuguese": "pt",
        "Russian": "ru",
        "Japanese": "ja",
        "Punjabi": "pa",
        "German": "de"
    }

    # defaults to chinese mapping if not found
    def get_language_code(language: str) -> str:
        for lang, code in top_languages.items():
            if language.lower() in lang.lower():
                return code
        return "zh-CN"
    
    return GoogleTranslator(source='auto', target=get_language_code(language)).translate(text)

def create_title_slide_translated(title_text: str, subtitle_text: str, prs, title_size: int, default_body_size: int = 8, language: str = 'chinese'):
    '''
    Creates a custom slide with a title text and body text.
    Usually used for song titles with small descriptions
    '''
    blank_slide = create_blank_slide(prs)

    title_text_translated = translate_text(title_text, language)
    subtitle_text_translated = translate_text(subtitle_text, language)

    subtitle_text = subtitle_text.replace("\n", " ")

    add_text_to_slide(blank_slide, f'{title_text}\n{title_text_translated}', prs, title_size, position_percent=0.2)
    add_text_to_slide(blank_slide, f'{subtitle_text}\n{subtitle_text_translated}', prs, default_body_size, position_percent=0.6)

    print(f'Adding: {title_text} | {title_text_translated}')

    return prs

def create_text_slide_translated(title_text, body_text, prs, title_size,
                                body_size, slide_number=0, total_slides=0, song_mode=False,
                                language="Chinese (Simplified)"):
    '''
    Creates a normal slide with just smaller body text and translation underneath.
    Result is similar to the "Content" slide but with bilingual support.
    Usually used in song lyrics slides.
    
    Args:
        title_text: English title text (not displayed but used for reference)
        body_text: English body text (lyrics)
        prs: PowerPoint presentation object
        title_size: Font size for title (reference only)
        body_size: Font size for body text
        slide_number: Current slide number (for navigation)
        total_slides: Total number of slides
        song_mode: Whether this is for song lyrics
        language: Target language for translation
    '''

    # Make a new blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Use the blank slide layout
    lyric_slide = prs.slides.add_slide(blank_slide_layout)

    # Calculate the center and top position for body text box
    body_width = prs.slide_width * 0.9
    body_height = prs.slide_height * 0.8
    body_left = (prs.slide_width - body_width) / 2
    body_top = prs.slide_height * 0.1  # Start from top with some margin

    # Add a text box for the body text
    body_box = lyric_slide.shapes.add_textbox(left=body_left, top=body_top, width=body_width, height=body_height)
    body_frame = body_box.text_frame
    
    # Process body text line by line for translation
    body_lines = body_text.strip().split('\n')
    translated_body_lines = []
    
    for line in body_lines:
        if line.strip():  # Only translate non-empty lines
            translated_line = translate_text(line, language)
            translated_body_lines.append(f"{line}\n{translated_line}")
        else:
            translated_body_lines.append(line)
    
    combined_body = '\n'.join(translated_body_lines)
    body_frame.text = combined_body

    # Enable text wrapping
    body_frame.word_wrap = True

    # Set font size and formatting for body text
    for i, paragraph in enumerate(body_frame.paragraphs):
        # Alternate between English (larger) and translated (smaller) text
        if i % 2 == 0:  # English lines
            paragraph.font.size = Pt(body_size)
            paragraph.font.bold = False
        else:  # Translated lines
            paragraph.font.size = Pt(body_size - 4)
            paragraph.font.bold = False
            paragraph.font.italic = True
        paragraph.alignment = PP_ALIGN.CENTER

    return prs

def append_song_to_powerpoint_translated(song_name, prs, title_size, font_size, max_lines=2, language="Chinese (Simplified)"):
    '''
    Will append a song's lyrics from a text file with all its lyrics to the current powerpoint with translations.
    This is a translated version of the original append_song_to_powerpoint function.
    Args:
        song_name: Name of the song
        prs: PowerPoint presentation object
        title_size: Font size for titles
        font_size: Font size for lyrics
        max_lines: Maximum lines per slide before splitting
        language: Target language for translation
    '''
    # ONLY FOR TRANSLATED VERSIONS - SOME THINGS ARE A BIT SMALL
    font_size = font_size * 1.2

    # Create a new song object with the required data
    try:
        new_song = song_object_from_name(song_name, 2)
    except FileNotFoundError:
        print(f"{song_name} doesn't seem to exist.")
        return

    # First slide of the song with title data and ccli data
    song_ccli_info = find_ccli(new_song.title.replace("\n", "").replace("(live)", "").strip().lower())
    prs = create_title_slide_translated(new_song.title.strip().lower().title().replace(' (Live)', ''), song_ccli_info, prs, title_size, 8, language)

    lyrics_clump = ''.join([line[1] for line in new_song.lyrics])
    translated_lyrics_text = ''

    # If this works - use this and if it does not we go back to using standard translate
    ai_translate = False
    
    try:
        translated_lyrics_text = translate_with_gemini(lyrics_clump, language)
        ai_translate = True
    except Exception as e:
        print(e)
        print("AI translation failed - using standard Google translate module")

    if ai_translate:
        print("Successfully translated with AI")
        
        # Parse the temporary file similar to append_song_to_powerpoint

        # should be [english line, translated line, english line, translated line etc]
        translated_lyrics_text_list = translated_lyrics_text.split('\n')

        # fuse 2 lines (english then translated) into 1 string for each slide

        slide_text_list = []
        for line_index in range(1, len(translated_lyrics_text_list), 2):
            concat_line = f'{translated_lyrics_text_list[line_index - 1].strip()}\n{translated_lyrics_text_list[line_index].strip()}'
            slide_text_list.append(concat_line)
        
        # Create slides from the translated lyrics
        for slide_number, lyrics in enumerate(slide_text_list):
            prs = create_text_slide(lyrics, lyrics, prs, title_size, font_size,
                                                slide_number=slide_number, total_slides=len(translated_lyrics_text_list),
                                                song_mode=True)
        
    else:
        # using old google translate method
        for slide_number, lyrics in enumerate(new_song.lyrics):
            prs = create_text_slide_translated(lyrics[0], lyrics[1], prs, title_size, font_size,
                                                slide_number=slide_number, total_slides=len(new_song.lyrics),
                                                song_mode=True, language=language)

    return prs